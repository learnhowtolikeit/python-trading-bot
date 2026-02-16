from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Create the Presentation
prs = Presentation()

# --- HELPER FUNCTION ---
def add_slide(title, content_items):
    slide_layout = prs.slide_layouts[1] # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.text = content_items[0]
    
    for item in content_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

# --- SLIDE 1: TITLE ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Project Tech Stack & Dependencies"
slide.placeholders[1].text = "A breakdown of the Python libraries powering the Trading Bot\nFebruary 2026"

# --- SLIDE 2: THE CONCEPT ---
add_slide("What is 'pip' and why do we use it?", [
    "The Concept: Python is a 'bare bones' engine.",
    "PIP (Pip Installs Packages): The App Store for Python.",
    "Why we use libraries: To avoid reinventing the wheel.",
    "Instead of writing complex math or API code from scratch, we download pre-built, tested tools."
])

# --- SLIDE 3: DATA LAYER ---
add_slide("1. The Data Feed: 'yfinance'", [
    "Command: pip install yfinance",
    "Purpose: The 'Eyes' of the bot.",
    "What it does: Connects to Yahoo Finance servers to fetch real-time price history.",
    "Context: Without this, we would need to manually enter stock prices or pay thousands for a Bloomberg terminal feed."
])

# --- SLIDE 4: ANALYTICS LAYER ---
add_slide("2. The Brain: 'pandas_ta'", [
    "Command: pip install pandas_ta",
    "Purpose: The 'Analyst' inside the bot.",
    "What it does: contains the math formulas for 130+ indicators (RSI, MACD, SMA).",
    "Context: It turns raw price numbers into actionable signals.",
    "Note: This library relies on 'pandas', the industry standard for data manipulation."
])

# --- SLIDE 5: NOTIFICATION LAYER ---
add_slide("3. The Voice: 'plyer'", [
    "Command: pip install plyer",
    "Purpose: The 'Messenger'.",
    "What it does: Sends commands to the Windows 10/11 Operating System to trigger a native desktop popup.",
    "Context: This allows the bot to run in the background (headless) while we work on other things."
])

# --- SLIDE 6: DOCUMENTATION LAYER ---
add_slide("4. The Scribe: 'python-pptx'", [
    "Command: pip install python-pptx",
    "Purpose: Automated Documentation.",
    "What it does: Allows Python to create, edit, and save PowerPoint files.",
    "Context: We use this to generate reports and lessons learned automatically, ensuring documentation never lags behind code."
])

# --- SLIDE 7: LESSONS LEARNED (DEPRECATED TOOLS) ---
add_slide("5. The Graveyard: 'win10toast'", [
    "Command: pip install win10toast",
    "Status: DEPRECATED (Removed).",
    "The Lesson: Not all libraries are maintained forever.",
    "The Conflict: It relied on 'pkg_resources', which Python 3.13 no longer supports.",
    "The Fix: We replaced it with 'plyer', teaching us the importance of modern, maintained dependencies."
])

# --- SAVE THE FILE ---
file_name = "Trading_Bot_Tech_Stack.pptx"
prs.save(file_name)
print(f"Success! Created {file_name}")