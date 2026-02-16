from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Create the Presentation
prs = Presentation()

# --- HELPER FUNCTION ---
def add_slide(title, content_items):
    slide_layout = prs.slide_layouts[1] # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.text = content_items[0]
    
    for item in content_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

# --- SLIDE 1: TITLE ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Python Trading Bot: Final Report"
slide.placeholders[1].text = "From Zero to Desktop Automation\nFebruary 2026"

# --- SLIDE 2: THE SETUP ---
add_slide("Phase 1: The Foundation", [
    "Goal: specific development environment setup.",
    "Installed Visual Studio Code & Python 3.13.",
    "Mastered 'pip install' to get external tools.",
    "Library Used: 'yfinance' (Yahoo Finance API)."
])

# --- SLIDE 3: FIRST STEPS ---
add_slide("Phase 2: The First Scripts", [
    "Script: 'price_check.py'",
    "Connected to the API to fetch live NVDA price.",
    "Challenge: The script only ran once and stopped.",
    "Lesson: Learned the difference between a Script vs. a Bot."
])

# --- SLIDE 4: AUTOMATION ---
add_slide("Phase 3: The 'Looping' Bot", [
    "Script: 'bot_loop.py'",
    "Added 'While True' loop to run 24/7.",
    "Added 'Time.Sleep' to manage processing power.",
    "Implemented 'Working Hours' logic (The Bouncer).",
    "Result: Bot runs autonomously on a schedule."
])

# --- SLIDE 5: INTELLIGENCE ---
add_slide("Phase 4: The Analyst Logic", [
    "Script: 'bot_crossover.py'",
    "Added Moving Average (SMA) calculation.",
    "Created 'Memory' logic: Only alert on Trend CHANGES.",
    "Implants: Learned to use 'False Memory' to test alerts.",
    "Result: Bot analyzes data, doesn't just read it."
])

# --- SLIDE 6: THE FINAL FIX (New!) ---
add_slide("Phase 5: Modern Alerts & Debugging", [
    "Script: 'bot_plyer_desktop_alert.py'",
    "Goal: Windows Desktop Notifications.",
    "Bug Encountered: 'win10toast' failed (Old library vs New Python).",
    "The Fix: Swapped to 'plyer' (Modern library).",
    "Result: Successful visual pop-up alerts on trend change."
])

# --- SAVE THE FILE ---
file_name = "Trading_Bot_Final_Presentation.pptx"
prs.save(file_name)
print(f"Success! Created {file_name}")