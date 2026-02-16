from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Create the Presentation
prs = Presentation()

# --- HELPER FUNCTION (To make adding slides easier) ---
def add_slide(title, content_items):
    # Use Layout 1 (Title + Content)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    slide.shapes.title.text = title
    
    # Add Content (Bullet Points)
    tf = slide.placeholders[1].text_frame
    tf.text = content_items[0] # First bullet
    
    for item in content_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

# --- SLIDE 1: TITLE SLIDE ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Python Trading Bot Project"
slide.placeholders[1].text = "Automating Analysis & Alerts\nCreated with Python"

# --- SLIDE 2: PHASE 1 - THE SETUP ---
add_slide("Phase 1: The Foundation", [
    "Goal: specific development environment setup.",
    "Installed Visual Studio Code (The Workshop).",
    "Installed Python 3 (The Engine).",
    "Learned 'pip install' to get tools (yfinance).",
    "Result: Successfully ran first Python command."
])

# --- SLIDE 3: PHASE 2 - FIRST STEPS ---
add_slide("Phase 2: Data Fetching", [
    "Script: 'price_check.py'",
    "Connected to Yahoo Finance API.",
    "Fetched live price for NVDA.",
    "Learned the difference between a Script (manual) and a Bot (automatic)."
])

# --- SLIDE 4: PHASE 3 - AUTOMATION ---
add_slide("Phase 3: The 'Looping' Bot", [
    "Script: 'bot_loop.py'",
    "Added 'While True' loop to run forever.",
    "Added 'Time.Sleep' to prevent crashing.",
    "Implemented 'Working Hours' (The Bouncer logic).",
    "Result: Bot watches the market 24/7 automatically."
])

# --- SLIDE 5: PHASE 4 - INTELLIGENCE ---
add_slide("Phase 4: Smart Alerts", [
    "Script: 'bot_desktop_alert.py'",
    "Added Moving Average (SMA) logic.",
    "Created 'Memory' so it only alerts on CHANGES.",
    "Integrated Windows Desktop Notifications.",
    "Result: Hands-free trading alerts."
])

# --- SAVE THE FILE ---
file_name = "Trading_Bot_Summary.pptx"
prs.save(file_name)
print(f"Success! Created {file_name} in your folder.")