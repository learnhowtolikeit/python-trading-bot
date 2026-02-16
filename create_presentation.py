from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 1. Create the Presentation Object
prs = Presentation()

# --- Helper Function to add a slide easily ---
def add_slide(title_text, content_text_list):
    # Layout 1 is "Title and Content" (standard bullet points)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the Title
    title = slide.shapes.title
    title.text = title_text
    
    # Set the Content (Bullet points)
    tf = slide.placeholders[1].text_frame
    tf.text = content_text_list[0] # First bullet
    
    # Add the rest of the bullets
    for item in content_text_list[1:]:
        p = tf.add_paragraph()
        p.text = item

# --- SLIDE 1: Title Slide ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Building My First Python Trading Bot"
slide.placeholders[1].text = "A Journey from Zero to Automation\nFebruary 2026"

# --- SLIDE 2: The Setup ---
add_slide("Phase 1: The Setup", [
    "Installed Visual Studio Code (The Workshop).",
    "Installed Python (The Engine).",
    "Learned to use the Terminal to install libraries.",
    "Key Library Used: 'yfinance' (Connects to Yahoo Finance API)."
])

# --- SLIDE 3: The First Script ---
add_slide("Phase 2: The First Script", [
    "Goal: Fetch a single stock price.",
    "Learned how to import libraries.",
    "Wrote code to 'Ask' Yahoo for NVDA data.",
    "Success: Printed '$182.81' to the terminal.",
    "Realization: This was a script, not a bot (it didn't loop)."
])

# --- SLIDE 4: The Automation ---
add_slide("Phase 3: The Loop (Bot Mode)", [
    "Goal: Make it run automatically.",
    "Added a 'While True' loop (Infinite Loop).",
    "Added 'Time.sleep(10)' to prevent crashing the computer.",
    "Result: The script now watches the market 24/7.",
    "It prints a new price update every 10 seconds."
])

# --- SLIDE 5: Next Steps ---
add_slide("Phase 4: Future Goals", [
    "Add 'Logic': Only notify me if price drops below $180.",
    "Paper Trading: Connect to a fake money account.",
    "Expansion: Move from Webull (Manual) to API Trading.",
    "Goal: Fully automated swing trading assistance."
])

# --- SAVE THE FILE ---
prs.save('My_Trading_Bot_Journey.pptx')
print("Success! Your PowerPoint file 'My_Trading_Bot_Journey.pptx' has been created.")