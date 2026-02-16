from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Create the Presentation
prs = Presentation()

# --- HELPER FUNCTION ---
def add_slide(title, content_items):
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.text = content_items[0]
    
    for item in content_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

# --- SLIDE 1: TITLE ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "The Developer's Journey"
slide.placeholders[1].text = "Building & Publishing a Python Trading Bot\nFrom Zero to GitHub"

# --- SLIDE 2: THE VISION ---
add_slide("Step 1: The Vision", [
    "Objective: Build an automated stock trading assistant.",
    "The Problem: We cannot watch the screen 24/7.",
    "The Solution: A Python script that watches, thinks, and alerts.",
    "Tools Selected: VS Code (Editor), Python 3.13 (Engine)."
])

# --- SLIDE 3: THE STRATEGY ---
add_slide("Step 2: The Logic (The 'Trifecta')", [
    "We programmed the bot to think like a disciplined trader.",
    "Indicator 1: SMA (20) - Is the trend up?",
    "Indicator 2: RSI (14) - Is the price fair?",
    "Indicator 3: MACD (12,26,9) - Is momentum building?",
    "Rule: The bot only signals when ALL three agree."
])

# --- SLIDE 4: THE AUTOMATION ---
add_slide("Step 3: Automation & Alerts", [
    "Challenge: Making it run forever.",
    "Solution: Implemented a 'While True' infinite loop.",
    "The Voice: Added 'plyer' to trigger Windows Desktop Notifications.",
    "The Result: The user can work on other tasks while the bot guards the portfolio."
])

# --- SLIDE 5: THE MEMORY ---
add_slide("Step 4: Persistent Logging", [
    "Challenge: How do we know what the bot did while we slept?",
    "Solution: Added a CSV Logger.",
    "Mechanism: Python opens 'trade_log.csv' and appends every Buy/Sell decision.",
    "Benefit: Creates a permanent audit trail for backtesting and review."
])

# --- SLIDE 6: VERSION CONTROL ---
add_slide("Step 5: Version Control (Git)", [
    "We turned our folder into a Repository.",
    "Git Init: Started tracking changes.",
    "Git Add/Commit: Saved 'snapshots' of our code.",
    "Why? To prevent data loss and allow us to 'rewind' if we break something."
])

# --- SLIDE 7: PUBLISHING ---
add_slide("Step 6: Publishing to the Cloud (GitHub)", [
    "Final Step: Uploading to the world.",
    "Remote Repo: Created a secure box on GitHub.com.",
    "Push: Sent our local code to the cloud.",
    "Authentication: Secured the connection using a Personal Access Token (PAT).",
    "Outcome: The code is now safe, shareable, and professional."
])

# --- SLIDE 8: NEXT STEPS ---
add_slide("Future Roadmap", [
    "Phase 2: Cloud Hosting (Running 24/7 on a server).",
    "Phase 3: Salesforce Integration (Logging trades as CRM records).",
    "Phase 4: Backtesting (Simulating performance on past data).",
    "Status: Phase 1 Complete."
])

# --- SAVE THE FILE ---
file_name = "My_Trading_Bot_Journey.pptx"
prs.save(file_name)
print(f"Success! Created {file_name}")