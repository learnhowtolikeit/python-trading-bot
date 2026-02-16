from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Create the Presentation
prs = Presentation()

# --- HELPER FUNCTION ---
def add_slide(title, content_items):
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.text = content_items[0]
    
    for item in content_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

# --- SLIDE 1: TITLE ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "The 'Trifecta' Trading Strategy"
slide.placeholders[1].text = "Combining SMA, RSI, and MACD for High-Probability Trades\nFebruary 2026"

# --- SLIDE 2: THE PROBLEM ---
add_slide("The Problem: Trading Blind", [
    "Most beginners only look at the Price.",
    "Risk 1: Buying a 'cheap' stock that is actually crashing (Catching a falling knife).",
    "Risk 2: Buying a 'strong' stock right before it reverses (FOMO).",
    "Solution: We need a 'Council of Elders'—3 distinct indicators to confirm every trade."
])

# --- SLIDE 3: INDICATOR 1 - SMA (The Map) ---
add_slide("1. SMA (Simple Moving Average)", [
    "Role: The Trend Filter.",
    "The Logic: Calculate the average price over the last 20 days.",
    "The Rule: If Price > SMA, the Trend is UP (Bullish).",
    "Why we need it: It keeps us on the right side of the market.",
    "Safety Check: We NEVER buy if the price is below the SMA."
])

# --- SLIDE 4: INDICATOR 2 - RSI (The Gas Gauge) ---
add_slide("2. RSI (Relative Strength Index)", [
    "Role: Momentum & Timing.",
    "The Logic: Measures how fast the price is moving (0-100 scale).",
    "RSI < 30 (Oversold): The stock is temporarily cheap (Good entry).",
    "RSI > 70 (Overbought): The stock is temporarily expensive (Good exit).",
    "Why we need it: It prevents us from buying at the very top."
])

# --- SLIDE 5: INDICATOR 3 - MACD (The Engine) ---
add_slide("3. MACD (Momentum Convergence Divergence)", [
    "Role: Momentum Confirmation.",
    "The Logic: Compares two moving averages to find shifts in power.",
    "The Signal: The 'Crossover'.",
    "   - MACD Line crosses ABOVE Signal Line = Bullish Momentum.",
    "   - MACD Line crosses BELOW Signal Line = Bearish Momentum.",
    "Why we need it: It confirms the move is real, not a fake-out."
])

# --- SLIDE 6: THE STRATEGY (CONFLUENCE) ---
add_slide("Putting It Together: The 'Trifecta' Setup", [
    "We ONLY trade when all three lights turn GREEN:",
    "1. Trend: Price is ABOVE the SMA (It's safe).",
    "2. Value: RSI is NOT Overbought (It's not too expensive).",
    "3. Momentum: MACD just Crossed UP (The move is starting).",
    "Result: Higher win rate, lower risk."
])

# --- SLIDE 7: PYTHON IMPLEMENTATION ---
add_slide("Automating the Strategy with Python", [
    "Library: 'pandas_ta' calculates all math instantly.",
    "The Bot's Logic:",
    "   IF (Price > SMA) AND (RSI < 70) AND (MACD_Cross == True):",
    "       SEND_ALERT('Golden Trade Opportunity')",
    "This removes human emotion and ensures discipline."
])

# --- SAVE THE FILE ---
file_name = "Bot_Strategy_Trifecta_Complete.pptx"
prs.save(file_name)
print(f"Success! Created {file_name}")