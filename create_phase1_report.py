from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Create the Presentation
prs = Presentation()

# --- HELPER FUNCTION ---
def add_slide(title, content_items):
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.text = content_items[0]
    
    for item in content_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

# --- SLIDE 1: TITLE ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Python Algorithmic Trading Bot"
slide.placeholders[1].text = "Phase 1 Completion Report: Desktop Automation\nFebruary 2026"

# --- SLIDE 2: THE OBJECTIVE ---
add_slide("Project Objective", [
    "Goal: specific development environment setup.",
    "Challenge: Create a bot that watches the stock market 24/7.",
    "Requirement 1: Must run autonomously (no human clicking).",
    "Requirement 2: Must analyze trends, not just price.",
    "Requirement 3: Must alert the user and log decisions."
])

# --- SLIDE 3: PHASE 1 - THE FOUNDATION ---
add_slide("Phase 1: The Build Environment", [
    "Tool: Visual Studio Code (The Workshop).",
    "Engine: Python 3.13.",
    "Package Manager: PIP (Pip Installs Packages).",
    "Key Library: 'yfinance' (Yahoo Finance API) to fetch live market data."
])

# --- SLIDE 4: PHASE 2 - THE INTELLIGENCE ---
add_slide("Phase 2: The 'Trifecta' Strategy", [
    "The Bot uses three indicators to make decisions (Confluence):",
    "1. SMA (Simple Moving Average): Determines the Trend direction.",
    "2. RSI (Relative Strength Index): Determines if price is cheap/expensive.",
    "3. MACD (Momentum): Confirms the strength of the move.",
    "Rule: ONLY trade when all three agree."
])

# --- SLIDE 5: PHASE 3 - THE VOICE ---
add_slide("Phase 3: Desktop Notifications", [
    "Challenge: The user cannot stare at the screen all day.",
    "Solution: Integrated Windows Notification Center.",
    "Library: 'plyer' (Replaced 'win10toast' due to dependency conflict).",
    "Result: Bot runs in the background and 'pops up' only when a trade is found."
])

# --- SLIDE 6: PHASE 4 - THE MEMORY (New!) ---
add_slide("Phase 4: Paper Trading Log", [
    "Feature: Automated CSV Logging.",
    "The Logic: Every 'Buy' or 'Sell' decision is written to 'trade_log.csv'.",
    "Data Captured: Timestamp, Ticker, Action, Price, Reasoning.",
    "Benefit: Allows for 'Paper Trading' (Simulated betting) to test strategy performance without financial risk."
])

# --- SLIDE 7: TECH STACK SUMMARY ---
add_slide("Final Tech Stack", [
    "Language: Python 3.13",
    "Data Source: yfinance",
    "Analytics: pandas & pandas_ta",
    "System Alerts: plyer",
    "Data Storage: CSV (Excel Compatible)",
    "Documentation: python-pptx"
])

# --- SAVE THE FILE ---
file_name = "Trading_Bot_Phase1_Complete.pptx"
prs.save(file_name)
print(f"Success! Created {file_name}")